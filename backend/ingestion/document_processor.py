import os
import re
import csv
import fitz  # PyMuPDF
import zipfile
import xml.etree.ElementTree as ET
from typing import List, Dict, Any
from dataclasses import dataclass, asdict

@dataclass
class DocumentChunk:
    text: str
    page: int
    chunk_index: int
    metadata: Dict[str, Any]

def extract_text_from_pdf(file_path: str) -> List[Dict[str, Any]]:
    """
    Extracts text page-by-page from a PDF using PyMuPDF.
    """
    pages_data = []
    try:
        doc = fitz.open(file_path)
        for page_idx, page in enumerate(doc):
            text = page.get_text()
            
            # Extract tables as structured text
            try:
                tables = page.find_tables()
                for table in tables:
                    table_text = []
                    for row in table.extract():
                        clean_row = [str(cell).strip().replace("\n", " ") for cell in row if cell is not None]
                        if any(clean_row):
                            table_text.append(" | ".join(clean_row))
                    if table_text:
                        text += "\n\n[TABLE DATA]\n" + "\n".join(table_text)
            except Exception:
                pass

            if not text.strip():
                text = f"[Image/Scanned Page - OCR Not Available for page {page_idx + 1}]"
            pages_data.append({
                "page": page_idx + 1,
                "text": text
            })
        doc.close()
    except Exception as e:
        raise RuntimeError(f"Error reading PDF file {file_path}: {e}")
    return pages_data

def extract_text_from_docx(file_path: str) -> List[Dict[str, Any]]:
    """
    Extracts text from a DOCX file using native zip/xml parsing.
    """
    try:
        with zipfile.ZipFile(file_path) as docx:
            xml_content = docx.read('word/document.xml')
            root = ET.fromstring(xml_content)

            ns = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main'}

            paragraphs = []
            for para in root.findall('.//w:p', ns):
                texts = [node.text for node in para.findall('.//w:t', ns) if node.text]
                if texts:
                    paragraphs.append("".join(texts))

            full_text = "\n".join(paragraphs)
            return [{"page": 1, "text": full_text}]
    except Exception as e:
        raise RuntimeError(f"Error reading DOCX file {file_path}: {e}")

def extract_text_from_txt(file_path: str) -> List[Dict[str, Any]]:
    """
    Extracts text from plain text file.
    """
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        return [{"page": 1, "text": text}]
    except Exception as e:
        raise RuntimeError(f"Error reading text file {file_path}: {e}")

def extract_text_from_csv(file_path: str) -> List[Dict[str, Any]]:
    """
    Reads a CSV file and converts each row to a natural-language sentence.
    This makes tabular data searchable via vector embeddings.
    E.g. "In row 5: Asset=P-101, Inspector=Rahul Kumar, Date=2024-03-12, Type=Inspection"
    """
    sentences = []
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            reader = csv.DictReader(f)
            for row_num, row in enumerate(reader, start=1):
                # Filter out empty cells
                fields = {k.strip(): v.strip() for k, v in row.items() if k and v and v.strip()}
                if not fields:
                    continue
                field_str = ", ".join(f"{k}={v}" for k, v in fields.items())
                sentences.append(f"Row {row_num}: {field_str}.")
        
        if not sentences:
            # Fallback: read as plain text
            return extract_text_from_txt(file_path)
        
        # Group sentences into page-sized chunks (100 rows per "page")
        pages = []
        chunk_size = 100
        for i in range(0, len(sentences), chunk_size):
            page_text = "\n".join(sentences[i:i + chunk_size])
            pages.append({"page": i // chunk_size + 1, "text": page_text})
        return pages
    except Exception as e:
        # If CSV parsing fails, fall back to plain text
        return extract_text_from_txt(file_path)

def extract_text_from_xlsx(file_path: str) -> List[Dict[str, Any]]:
    """
    Reads an Excel file using openpyxl and converts each sheet's rows 
    to natural-language sentences for semantic search.
    """
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        pages = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Extract header row
            headers = []
            rows_data = []
            
            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                if all(cell is None for cell in row):
                    continue  # skip empty rows
                
                if not headers:
                    # First non-empty row = headers
                    headers = [str(c).strip() if c is not None else f"Col{i}" for i, c in enumerate(row)]
                    continue
                
                # Convert row to natural language
                fields = {}
                for i, cell in enumerate(row):
                    if i < len(headers) and cell is not None:
                        val = str(cell).strip()
                        if val and val.lower() not in ("none", "null", ""):
                            fields[headers[i]] = val
                
                if fields:
                    field_str = ", ".join(f"{k}={v}" for k, v in fields.items())
                    rows_data.append(f"In sheet '{sheet_name}', row {row_idx}: {field_str}.")
            
            if rows_data:
                # Group into page-sized chunks
                chunk_size = 100
                for i in range(0, len(rows_data), chunk_size):
                    page_text = "\n".join(rows_data[i:i + chunk_size])
                    pages.append({"page": len(pages) + 1, "text": f"Sheet: {sheet_name}\n\n{page_text}"})
        
        if not pages:
            return [{"page": 1, "text": f"[Excel file {os.path.basename(file_path)} appears to be empty]"}]
        return pages
    except ImportError:
        raise RuntimeError("openpyxl is required for XLSX support. Install it with: pip install openpyxl")
    except Exception as e:
        raise RuntimeError(f"Error reading XLSX file {file_path}: {e}")

def extract_text_from_pptx(file_path: str) -> List[Dict[str, Any]]:
    """
    Reads a PowerPoint file using python-pptx and extracts text from each slide.
    Each slide becomes one "page" with title + body + speaker notes.
    """
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        pages = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            
            # Extract title
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_texts.append(f"Slide Title: {slide.shapes.title.text.strip()}")
            
            # Extract body text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if shape.shape_id != (slide.shapes.title.shape_id if slide.shapes.title else -1):
                        slide_texts.append(text)
            
            # Extract speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_texts.append(f"Speaker Notes: {notes_text}")
            
            if slide_texts:
                pages.append({
                    "page": slide_num,
                    "text": "\n".join(slide_texts)
                })
        
        if not pages:
            return [{"page": 1, "text": f"[PowerPoint file {os.path.basename(file_path)} appears to have no text content]"}]
        return pages
    except ImportError:
        raise RuntimeError("python-pptx is required for PPTX support. Install it with: pip install python-pptx")
    except Exception as e:
        raise RuntimeError(f"Error reading PPTX file {file_path}: {e}")

def detect_industrial_sections(text: str) -> List[Dict[str, Any]]:
    """
    Split text into sections using heading patterns.
    Look for patterns like:
    - ALL CAPS lines (typical SOP headings)
    - Numbered sections: "1.", "1.1", "Section 1"
    - Warning/Caution blocks: lines starting with WARNING:, CAUTION:, NOTE:
    - Work Order entries: "WO-", "Work Order", "Job:"
    - Procedure steps: lines starting with digits followed by period
    """
    HEADING_PATTERNS = [
        r'^\d+\.\d*\s+[A-Z]',           # "1.2 PROCEDURE TITLE"
        r'^[A-Z][A-Z\s]{4,}$',           # "EMERGENCY SHUTDOWN PROCEDURE"
        r'^(WARNING|CAUTION|NOTE)\s*:',  # safety blocks
        r'^(WO-|Work Order|Job:)',        # work orders
        r'^Section\s+\d+',               # "Section 3"
    ]

    RISK_SIGNALS = [
        "catastrophic", "emergency", "critical failure", "immediate action",
        "hazardous", "do not operate", "lockout", "LOTO",
        "[CONTENT MISSING]", "[TBD]", "pending safety review",
        "draft document", "awaiting input"
    ]

    combined_pattern = "|".join(HEADING_PATTERNS)
    lines = text.split('\n')
    sections = []
    
    current_heading = None
    current_text = []

    def _determine_section_type(heading, text_block):
        content = (str(heading) + " " + text_block).lower()
        if re.search(r'^(warning|caution|note)\s*:', str(heading), re.IGNORECASE):
            return "warning"
        if re.search(r'^(wo-|work order|job:)', str(heading), re.IGNORECASE):
            return "work_order"
        if "procedure" in content or re.search(r'^\d+\.\d*', str(heading)):
            return "procedure"
        if "maintenance" in content or "inspection" in content:
            return "maintenance"
        if "compliance" in content or "regulation" in content or "iso" in content or "osha" in content:
            return "compliance"
        return "general"

    def _flush_section():
        if current_text:
            text_block = "\n".join(current_text)
            sec_type = _determine_section_type(current_heading, text_block)
            has_risk = any(signal.lower() in text_block.lower() for signal in RISK_SIGNALS)
            sections.append({
                "heading": current_heading,
                "text": text_block,
                "section_type": sec_type,
                "has_risk_signal": has_risk
            })

    for line in lines:
        if re.match(combined_pattern, line.strip(), re.IGNORECASE):
            _flush_section()
            current_heading = line.strip()
            current_text = [line]
        else:
            current_text.append(line)
            
    _flush_section()
    
    if not sections:
        sections.append({
            "heading": None,
            "text": text,
            "section_type": "general",
            "has_risk_signal": any(signal.lower() in text.lower() for signal in RISK_SIGNALS)
        })
        
    return sections


def semantic_chunk_section(section: Dict[str, Any], max_words: int = 200, overlap_words: int = 30) -> List[Dict[str, Any]]:
    """
    Chunk a single section with word-count limit and overlap.
    """
    text = section["text"]
    heading = section.get("heading")
    paragraphs = re.split(r'\n\n+', text)
    
    chunks = []
    current_chunk_sentences = []
    current_word_count = 0
    
    for para in paragraphs:
        para = para.strip()
        if not para: continue
        
        sentences = re.split(r'(?<=\.)\s+(?=[A-Z])', para)
        
        for sentence in sentences:
            sentence_words = sentence.split()
            sentence_word_count = len(sentence_words)
            
            if current_word_count + sentence_word_count > max_words and current_chunk_sentences:
                chunk_text = " ".join(current_chunk_sentences)
                if heading:
                    chunk_text = f"[{heading}]\n" + chunk_text
                
                # Minimum chunk length filter (5 words)
                if len(chunk_text.split()) >= 5:
                    chunks.append({
                        "text": chunk_text,
                        "heading": heading,
                        "section_type": section["section_type"],
                        "has_risk_signal": section["has_risk_signal"]
                    })
                
                # Backtrack for overlap (keep sentences until we have at least overlap_words)
                overlap_sentences = []
                overlap_count = 0
                for s in reversed(current_chunk_sentences):
                    overlap_sentences.insert(0, s)
                    overlap_count += len(s.split())
                    if overlap_count >= overlap_words:
                        break
                        
                current_chunk_sentences = overlap_sentences + [sentence]
                current_word_count = overlap_count + sentence_word_count
            else:
                current_chunk_sentences.append(sentence)
                current_word_count += sentence_word_count
                
    if current_chunk_sentences:
        chunk_text = " ".join(current_chunk_sentences)
        if heading:
            chunk_text = f"[{heading}]\n" + chunk_text
            
        if len(chunk_text.split()) >= 5:
            chunks.append({
                "text": chunk_text,
                "heading": heading,
                "section_type": section["section_type"],
                "has_risk_signal": section["has_risk_signal"]
            })
        
    return chunks


def chunk_text(text: str, page_number: int, start_idx: int = 0, chunk_size: int = 512, overlap: int = 64) -> List[DocumentChunk]:
    """
    Splits text into structure-aware semantic chunks.
    """
    sections = detect_industrial_sections(text)
    
    chunks = []
    current_idx = start_idx
    
    for section in sections:
        section_chunks = semantic_chunk_section(section, max_words=chunk_size, overlap_words=overlap)
        for s_chunk in section_chunks:
            chunks.append(DocumentChunk(
                text=s_chunk["text"],
                page=page_number,
                chunk_index=current_idx,
                metadata={
                    "heading": s_chunk.get("heading"),
                    "section_type": s_chunk["section_type"],
                    "has_risk_signal": s_chunk["has_risk_signal"]
                }
            ))
            current_idx += 1
            
    return chunks

def process_file(file_path: str, filename: str) -> List[Dict[str, Any]]:
    """
    Identifies the file type, extracts text, chunks it, and returns the chunk list.
    Supports: PDF, DOCX, TXT, LOG, CSV, XLSX, XLS, PPTX, PPT.
    """
    ext = os.path.splitext(filename)[1].lower()

    if ext == ".pdf":
        pages = extract_text_from_pdf(file_path)
    elif ext == ".docx":
        pages = extract_text_from_docx(file_path)
    elif ext in (".txt", ".log", ".json"):
        pages = extract_text_from_txt(file_path)
    elif ext == ".csv":
        pages = extract_text_from_csv(file_path)
    elif ext in (".xlsx", ".xls"):
        pages = extract_text_from_xlsx(file_path)
    elif ext in (".pptx", ".ppt"):
        pages = extract_text_from_pptx(file_path)
    else:
        # Fallback to text reading for unknown formats
        pages = extract_text_from_txt(file_path)

    all_chunks = []
    chunk_index = 0

    for page_data in pages:
        page_text = page_data["text"]
        page_num = page_data["page"]

        if not page_text.strip():
            continue

        page_chunks = chunk_text(page_text, page_num, start_idx=chunk_index)
        for chunk in page_chunks:
            chunk_dict = asdict(chunk)
            chunk_dict["metadata"].update({
                "source_file": filename,
                "file_type": ext
            })
            all_chunks.append(chunk_dict)
            chunk_index += 1

    return all_chunks

def process_document(file_path: str, document_id: Any) -> List[Dict[str, Any]]:
    """
    Compatibility wrapper around process_file that normalizes keys.
    """
    filename = os.path.basename(file_path)
    chunks = process_file(file_path, filename)
    normalized_chunks = []
    for chunk in chunks:
        normalized_chunks.append({
            "text": chunk["text"],
            "page_number": chunk.get("page", 1),
            "chunk_index": chunk["chunk_index"],
            "metadata": chunk.get("metadata", {})
        })
    return normalized_chunks
